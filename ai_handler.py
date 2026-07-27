import os
import requests
from docx import Document
from pptx import Presentation
from pptx.util import Inches

from config import (
    AI_MODEL,
    OPENAI_API_KEY,
    CLAUDE_API_KEY,
    GEMINI_API_KEY,
)

# Pexels kaliti to'g'ridan-to'g'ri shu yerda
PEXELS_API_KEY = "EwSRENDIIVaujdEYjtp5WNAr26n67yI5GIJF7oK8gUR1b1yln3z3uSw3"

# ==========================================================
#  QO'LLAB-QUVVATLANADIGAN TILLAR
# ==========================================================

LANGUAGE_NAMES = {
    "uz": "o'zbek",
    "ru": "rus (Русский)",
    "en": "ingliz (English)",
    "kaa": "qoraqalpoq (Qaraqalpaqsha)",
    "kg": "qirg'iz (Кыргызcha)",
    "kk": "qozoq (Қазақша)",
    "tg": "tojik (Тоҷикӣ)",
}


def get_language_name(lang_code):
    return LANGUAGE_NAMES.get(lang_code, "o'zbek")

# ==========================================================
#  UCH XIL AI PROVAYDER: OpenAI -> Claude -> Gemini (fallback)
#  Birinchisi ishlamasa (masalan kredit tugasa), avtomatik
#  keyingisiga o'tadi.
# ==========================================================

def _try_openai(system_prompt, prompt):
    from openai import OpenAI
    client = OpenAI(api_key=OPENAI_API_KEY)
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": prompt}
    ]
    response = client.chat.completions.create(model=AI_MODEL, messages=messages)
    return response.choices[0].message.content


def _try_claude(system_prompt, prompt):
    # To'g'ridan-to'g'ri REST API orqali (kutubxona bilan bog'liq muammolarni oldini olish uchun)
    headers = {
        "x-api-key": CLAUDE_API_KEY,
        "anthropic-version": "2023-06-01",
        "content-type": "application/json",
    }
    payload = {
        "model": "claude-sonnet-4-6",
        "max_tokens": 8192,
        "system": system_prompt,
        "messages": [{"role": "user", "content": prompt}]
    }
    resp = requests.post(
        "https://api.anthropic.com/v1/messages",
        headers=headers, json=payload, timeout=60
    )
    data = resp.json()
    if "content" in data:
        return data["content"][0]["text"]
    raise Exception(data.get("error", {}).get("message", str(data)))


def _try_gemini(system_prompt, prompt):
    # To'g'ridan-to'g'ri REST API orqali (Windows'dagi DLL muammosini oldini olish uchun)
    url = (
        "https://generativelanguage.googleapis.com/v1beta/models/"
        f"gemini-2.0-flash:generateContent?key={GEMINI_API_KEY}"
    )
    payload = {
        "system_instruction": {"parts": [{"text": system_prompt}]},
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {"maxOutputTokens": 8192}
    }
    resp = requests.post(url, json=payload, timeout=60)
    data = resp.json()
    if "candidates" in data:
        return data["candidates"][0]["content"]["parts"][0]["text"]
    raise Exception(data.get("error", {}).get("message", str(data)))


# ==========================================================
#  HAR BIR HUJJAT TURI UCHUN AKADEMIK YOZISH TARTIBI (ANDOZA)
# ==========================================================

STRUCTURE_GUIDES = {
    "Dars ishlanmasi": """
Quyidagi tuzilishda, to'liq dars ishlanmasi (konspekt) yozing:
1. Mavzu
2. Darsning maqsadlari
3. Dars turi va metodi
4. Darsning borishi (tashkiliy qism, mavzu bayoni, mustahkamlash, uyga vazifa, baholash)
Har bir bo'limni aniq sarlavha bilan ajratib yozing.
""",

    "Dars ishlanmasi - Ma'ruza": """
Ma'ruza mashg'ulotini AYNAN quyidagi tuzilishda, sarlavhalarni katta harflar bilan
ajratib yozing. Har bir qismdagi so'z soni talabiga QAT'IY AMAL QILING:

MA'RUZA MASHG'ULOTI

MAVZU: (mavzu nomini yozing)

MAQSAD: (AYNAN 22 TA SO'ZDAN IBORAT matn)

VAZIFALARI: (aniq va lo'nda vazifalar, raqamlangan ro'yxat ko'rinishida)

REJA:
1. (reja bandi nomi)
2. (reja bandi nomi)
3. (reja bandi nomi)
4. (reja bandi nomi)

TAYANCH IBORALAR: (AYNAN 30 TA SO'ZDAN IBORAT, mavzuga oid asosiy atama va iboralar
ro'yxati matn ko'rinishida)

Endi REJA BANDLARI BO'YICHA MATN (har bir band alohida sarlavha bilan, reja nomini
qatorda ko'rsating):
1-reja bandi nomi bo'yicha: AYNAN 180 TA SO'ZDAN IBORAT matn, chuqur va tushunarli
2-reja bandi nomi bo'yicha: AYNAN 180 TA SO'ZDAN IBORAT matn, chuqur va tushunarli
3-reja bandi nomi bo'yicha: AYNAN 180 TA SO'ZDAN IBORAT matn, chuqur va tushunarli
4-reja bandi nomi bo'yicha: AYNAN 180 TA SO'ZDAN IBORAT matn, chuqur va tushunarli

NAZORAT UCHUN SAVOLLAR: (aynan 10 ta savol, raqamlangan)

AMALIY TOPSHIRIQ (KEYS-STADI): (mavzuga oid real hayotiy vaziyat/keys namunasini yozing)

MAVZU YUZASIDAN TEST: (aynan 10 ta test savoli, har biri 4 ta variant (A, B, C, D)
bilan, har bir savoldan keyin to'g'ri javobni ham ko'rsating)

MUSTAQIL IZLANISHLAR UCHUN MAVZULAR: (aynan 2 ta qo'shimcha mavzu, raqamlangan)

TAVSIYA ETILGAN ADABIYOTLAR: (aynan 4 yoki 5 ta manba — albatta mavzuga oid MAHALLIY
(o'zbekistonlik) OLIMLAR tomonidan yaratilgan adabiyotlar bo'lsin, standart bibliografik
shaklda, raqamlangan)
""",

    "Dars ishlanmasi - Amaliy mashg'ulot": """
Amaliy mashg'ulotni AYNAN quyidagi tuzilishda, sarlavhalarni katta harflar bilan
ajratib yozing. Har bir qismdagi so'z soni talabiga QAT'IY AMAL QILING:

AMALIY MASHG'ULOT

MAVZU: (mavzu nomi)

Mashg'ulotning maqsadi: (AYNAN 50 TA SO'ZDAN IBORAT matn)

Topshiriqlar: (aynan 3 ta aniq topshiriq, raqamlangan: 1, 2, 3)

Mashg'ulotning borishi:
1-bosqich: (AYNAN 50 TA SO'ZDAN IBORAT matn)
2-bosqich: (AYNAN 25 TA SO'ZDAN IBORAT matn)
3-bosqich: (AYNAN 25 TA SO'ZDAN IBORAT matn)

AMALIY TOPSHIRIQLAR VA METODIK O'YINLAR
Topshiriq: (mavzuga oid sport va harakatli o'yinlarni tavsiflang — kamida 2 ta o'yin,
har birining qisqacha qoidasi va mavzu bilan bog'liqligini tushuntiring)

Yakuniy xulosa: (AYNAN 25 TA SO'ZDAN IBORAT matn)

Nazorat savollari: (aynan 6 ta savol, raqamlangan)

TAVSIYA ETILGAN ADABIYOTLAR: (aynan 4 yoki 5 ta manba — albatta mavzuga oid MAHALLIY
(o'zbekistonlik) OLIMLAR tomonidan yaratilgan adabiyotlar bo'lsin, raqamlangan)

Ilova: (mavzuga oid rasm yoki texnologik xaritaning tavsifini matn shaklida tasvirlab
bering)
""",

    "Dars ishlanmasi - Seminar": """
Seminar mashg'ulotini AYNAN quyidagi tuzilishda, sarlavhalarni katta harflar bilan
ajratib yozing. Har bir qismdagi so'z soni talabiga QAT'IY AMAL QILING:

SEMINAR MASHG'ULOTI

MAVZU: (mavzu nomi)

Reja:
1. (band nomi)
2. (band nomi)
3. (band nomi)
4. (band nomi)

Mashg'ulotning maqsadi: (AYNAN 60 TA SO'ZDAN IBORAT matn)

Mashg'ulot turi: (masalan: interfaol seminar, muammoli seminar va h.k.)

Mashg'ulot muammosi: (AYNAN 30 TA SO'ZDAN IBORAT matn)

O'quv maqsadlari: (2-3 band, raqamlangan)

Seminar reja: (seminar davomidagi asosiy bosqichlar, qisqa ro'yxat ko'rinishida)

Asosiy tushunchalar: (aynan 5 ta atama, har biri qisqa ta'rifi bilan)

Guruh topshiriqlari: (kamida 3 ta guruh uchun, har biriga alohida topshiriq)

Seminar metodlari: (qo'llaniladigan metodlar ro'yxati va qisqa tavsifi)

Uyga vazifa: (mavzuga asoslangan aniq topshiriq)

Mini-insho yozish (1-1.5 sahifa): (insho uchun aniq mavzu/savol beriladi)

Seminarlarda muhokama qilish uchun qo'shimcha savollar: (aynan 5 ta savol, raqamlangan)

TAVSIYA ETILGAN ADABIYOTLAR: (aynan 4 yoki 5 ta manba — albatta mavzuga oid MAHALLIY
(o'zbekistonlik) OLIMLAR tomonidan yaratilgan adabiyotlar bo'lsin, raqamlangan)
""",

    "Dars ishlanmasi - Laboratoriya mashg'uloti": """
Laboratoriya mashg'ulotini AYNAN quyidagi tuzilishda, sarlavhalarni katta harflar
bilan ajratib yozing. Har bir qismdagi so'z soni talabiga QAT'IY AMAL QILING:

LABORATORIYA MASHG'ULOTI

MAVZU: (mavzu nomi)

Reja:
1. (band nomi)
2. (band nomi)
3. (band nomi)
4. (band nomi)

Mashg'ulotning maqsadi: (AYNAN 60 TA SO'ZDAN IBORAT matn)

Mashg'ulot turi: (masalan: tajriba-amaliy laboratoriya mashg'uloti va h.k.)

Mashg'ulot muammosi: (AYNAN 30 TA SO'ZDAN IBORAT matn)

O'quv maqsadlari: (2-3 band, raqamlangan)

Laboratoriya reja: (mashg'ulot davomidagi asosiy bosqichlar, qisqa ro'yxat ko'rinishida)

Asosiy tushunchalar: (aynan 5 ta atama, har biri qisqa ta'rifi bilan)

Guruh topshiriqlari: (kamida 3 ta guruh uchun, har biriga alohida topshiriq)

Laboratoriya metodlari: (qo'llaniladigan metod va uskunalar ro'yxati va qisqa tavsifi)

Uyga vazifa: (mavzuga asoslangan aniq topshiriq)

Mini-insho yozish (1-1.5 sahifa): (insho uchun aniq mavzu/savol beriladi)

Laboratoriyalarda muhokama qilish uchun qo'shimcha savollar: (aynan 5 ta savol,
raqamlangan)

TAVSIYA ETILGAN ADABIYOTLAR: (aynan 4 yoki 5 ta manba — albatta mavzuga oid MAHALLIY
(o'zbekistonlik) OLIMLAR tomonidan yaratilgan adabiyotlar bo'lsin, raqamlangan)
""",

    "Mavzu bo'yicha slayd - Ma'ruza": """
Ma'ruza mavzusi bo'yicha TAQDIMOT (PowerPoint) tarkibini tayyorlang. Natija albatta
15-20 TA ALOHIDA SLAYDGA mos bo'lingan bo'lishi kerak. Har bir slayd quyidagi formatda,
bir-biridan BO'SH QATOR bilan ajratilgan holda berilsin:

SLAYD SARLAVHASI
- Qisqa va lo'nda band (to'liq gap emas, taqdimotga mos qisqa ibora)
- Qisqa va lo'nda band
- Qisqa va lo'nda band
(har bir slaydda 3-5 ta qisqa band yetarli)

Slaydlar ketma-ketligi (jami 15-20 ta slaydga yetguncha davom eting):
1. Sarlavha slaydi (mavzu nomi + "Ma'ruza")
2. Maqsad va vazifalar
3. Reja (asosiy bandlar ro'yxati)
4. Tayanch iboralar
5-12. Reja bandlari bo'yicha asosiy mazmun (har bir band uchun kamida 2 ta slayd,
   band ichidagi ma'lumotni bir necha slaydga bo'lib bering)
13. Amaliy topshiriq / KEYS-STADI namunasi
14. Nazorat savollari (qisqa ro'yxat)
15. Mustaqil izlanish uchun mavzular
16. Tavsiya etilgan adabiyotlar
Zarurat bo'lsa slaydlar sonini 20 tagacha oshiring, lekin 15 tadan kam bo'lmasin.

MUHIM: Agar biror slaydda jadval ko'rinishida taqdim etish mumkin bo'lgan ma'lumot
bo'lsa (masalan reja, taqqoslash, ko'rsatkichlar), uni albatta jadval formatida
("ustun1 | ustun2 | ustun3" ko'rinishida, birinchi qator sarlavha bo'lsin) bering.
""",

    "Mavzu bo'yicha slayd - Amaliy mashg'ulot": """
Amaliy mashg'ulot mavzusi bo'yicha TAQDIMOT (PowerPoint) tarkibini tayyorlang. Natija
albatta 12-15 TA ALOHIDA SLAYDGA mos bo'lingan bo'lishi kerak. Har bir slayd quyidagi
formatda, bir-biridan BO'SH QATOR bilan ajratilgan holda berilsin:

SLAYD SARLAVHASI
- Qisqa va lo'nda band
- Qisqa va lo'nda band
- Qisqa va lo'nda band

Slaydlar ketma-ketligi (jami 12-15 ta slaydga yetguncha davom eting):
1. Sarlavha slaydi (mavzu nomi + "Amaliy mashg'ulot")
2. Mashg'ulotning maqsadi
3. Topshiriqlar
4-7. Mashg'ulotning borishi (1-4 bosqichlar, har biri alohida slaydda)
8. Amaliy topshiriqlar va metodik o'yinlar
9. Blits-so'rov
10. Yakuniy xulosa
11-12. Nazorat savollari
13. Tavsiya etilgan adabiyotlar
Zarurat bo'lsa slaydlar sonini 15 tagacha oshiring, lekin 12 tadan kam bo'lmasin.

MUHIM: Agar biror slaydda jadval ko'rinishida taqdim etish mumkin bo'lgan ma'lumot
bo'lsa, uni albatta jadval formatida ("ustun1 | ustun2 | ustun3" ko'rinishida,
birinchi qator sarlavha bo'lsin) bering.
""",

    "Mavzu bo'yicha slayd - Seminar": """
Seminar mavzusi bo'yicha TAQDIMOT (PowerPoint) tarkibini tayyorlang. Natija albatta
12-15 TA ALOHIDA SLAYDGA mos bo'lingan bo'lishi kerak. Har bir slayd quyidagi formatda,
bir-biridan BO'SH QATOR bilan ajratilgan holda berilsin:

SLAYD SARLAVHASI
- Qisqa va lo'nda band
- Qisqa va lo'nda band
- Qisqa va lo'nda band

Slaydlar ketma-ketligi (jami 12-15 ta slaydga yetguncha davom eting):
1. Sarlavha slaydi (mavzu nomi + "Seminar")
2. Reja
3. Mashg'ulotning maqsadi, turi va muammosi
4. Seminar rejasi (bosqichlar va vaqt taqsimoti)
5. Asosiy tushunchalar
6-9. Guruh topshiriqlari (har bir guruh uchun alohida slayd)
10. Seminar metodlari
11. Baholash mezonlari
12. Uyga vazifa
13. Qo'shimcha muhokama savollari
14. Tavsiya etilgan adabiyotlar
Zarurat bo'lsa slaydlar sonini 15 tagacha oshiring, lekin 12 tadan kam bo'lmasin.

MUHIM: "Seminar rejasi" va "Baholash mezonlari" slaydlarini albatta jadval formatida
("ustun1 | ustun2 | ustun3" ko'rinishida, birinchi qator sarlavha bo'lsin) bering.
""",

    "Mavzu bo'yicha slayd - Laboratoriya mashg'uloti": """
Laboratoriya mashg'uloti mavzusi bo'yicha TAQDIMOT (PowerPoint) tarkibini tayyorlang.
Natija albatta 12-15 TA ALOHIDA SLAYDGA mos bo'lingan bo'lishi kerak. Har bir slayd
quyidagi formatda, bir-biridan BO'SH QATOR bilan ajratilgan holda berilsin:

SLAYD SARLAVHASI
- Qisqa va lo'nda band
- Qisqa va lo'nda band
- Qisqa va lo'nda band

Slaydlar ketma-ketligi (jami 12-15 ta slaydga yetguncha davom eting):
1. Sarlavha slaydi (mavzu nomi + "Laboratoriya mashg'uloti")
2. Mashg'ulotning maqsadi
3. Kerakli asbob-uskunalar va materiallar
4. Xavfsizlik qoidalari
5. Nazariy qism (qisqacha)
6-10. Ish tartibi (bosqichma-bosqich, har biri alohida slaydda)
11. Natijalarni qayd etish jadvali
12. Xulosa chiqarish uchun savollar
13. Nazorat savollari
14. Tavsiya etilgan adabiyotlar
Zarurat bo'lsa slaydlar sonini 15 tagacha oshiring, lekin 12 tadan kam bo'lmasin.

MUHIM: "Natijalarni qayd etish jadvali" slaydini albatta jadval formatida
("ustun1 | ustun2 | ustun3" ko'rinishida, birinchi qator sarlavha bo'lsin) bering.
""",

    "Maqola": """
Ilmiy maqolani xalqaro IMRAD (Introduction, Methods, Results, And Discussion) talablariga
mos, quyidagi aniq tuzilishda yozing:

UDK (УДК) — mavzuga mos taxminiy UDK indeksini ko'rsating (masalan: UDK 004.8)

SARLAVHA — uch tilda bering:
  1) (tanlangan tilda)
  2) Rus tilida
  3) Ingliz (English) tilida

MUALLIF: [Muallif F.I.Sh. shu yerga kiritiladi]
MUASSASA: [Muassasa nomi shu yerga kiritiladi]

ANNOTATSIYA — uch tilda, har biri 150-200 so'zdan iborat:
  1) (tanlangan tilda annotatsiya)
  2) Аннотация (rus tilida)
  3) Abstract (ingliz tilida)

KALIT SO'ZLAR — uch tilda, har birida 5-7 ta so'z:
  1) Kalit so'zlar: (tanlangan tilda)
  2) Ключевые слова: (rus tilida)
  3) Keywords: (ingliz tilida)

KIRISH (Introduction) — muammoning dolzarbligi, ilmiy bo'shliq, maqsad
MATERIALLAR VA METODLAR (Methods) — tadqiqot metodologiyasi, ma'lumotlar manbai
NATIJALAR (Results) — asosiy topilmalar, tahlil natijalari
MUHOKAMA (Discussion) — natijalarning talqini, boshqa tadqiqotlar bilan solishtirish
XULOSA (Conclusion) — asosiy xulosalar va tavsiyalar
FOYDALANILGAN ADABIYOTLAR RO'YXATI (kamida 8-10 ta manba, xalqaro standart shaklda)
""",

    "Tezis": """
Qisqa ilmiy tezis uchun quyidagi tuzilishda yozing:
1. Sarlavha
2. Kirish (muammoning dolzarbligi)
3. Asosiy g'oyalar (aniq va lo'nda, punktlar bilan)
4. Xulosa
5. Foydalanilgan adabiyotlar (3-5 ta manba)
""",

    "Mustaqil ish": """
Mustaqil ish uchun quyidagi tuzilishda yozing:
1. Kirish (mavzuning ahamiyati, ishning maqsadi)
2. Asosiy qism (mavzuni bo'limlarga bo'lib, chuqur va tushunarli tarzda yoritish)
3. Xulosa (asosiy fikrlarning qisqacha yakuni)
4. Foydalanilgan adabiyotlar ro'yxati
""",

    "Kurs ishi": """
Kurs ishi uchun quyidagi to'liq akademik tuzilishda yozing (hajmi 30-35 betga mo'ljallangan,
shuning uchun har bir qismni imkon qadar batafsil yoriting):
1. Mundarija
2. Kirish (mavzuning dolzarbligi, ishning maqsadi va vazifalari, tadqiqot ob'ekti/predmeti)
3. I BOB — Nazariy qism (mavzu bo'yicha ilmiy adabiyotlar tahlili, asosiy tushunchalar)
4. II BOB — Amaliy/tahliliy qism (masala yechimi, misollar, tahlil, statistik yoki amaliy ma'lumotlar)
5. Xulosa va takliflar
6. Foydalanilgan adabiyotlar ro'yxati (kamida 8-10 ta manba)
7. Ilovalar (agar kerak bo'lsa, qaysi turdagi ilova kerakligini ko'rsating)
""",

    "Bitiruv malakaviy ishi": """
Bitiruv malakaviy ishi (BMI) uchun quyidagi rasmiy OAK talablariga mos, TO'LIQ va BATAFSIL
akademik tuzilishda yozing (umumiy hajmi 70-80 betga mo'ljallangan, shuning uchun har bir
kichik bo'limni (masalan 1.1, 1.2) alohida, chuqur va keng yoriting):

MUNDARIJA (I, II, III bob asosida)

KIRISH
- Mavzuning dolzarbligi va zarurati
- Tadqiqotning maqsadi va vazifalari
- Tadqiqot ob'ekti va predmeti
- Ishning ilmiy-amaliy ahamiyati

I BOB (nazariy-metodologik asoslar)
  1.1. (kichik mavzu — mavzuning nazariy asoslari)
  1.2. (kichik mavzu — mavzu bo'yicha ilmiy qarashlar tahlili)
  1.3. (kichik mavzu — xorijiy va mahalliy tajriba)
  I bob bo'yicha xulosa

II BOB (amaliy tahlil)
  2.1. (kichik mavzu — hozirgi holat tahlili)
  2.2. (kichik mavzu — statistik/amaliy ma'lumotlar tahlili)
  2.3. (kichik mavzu — muammolarni aniqlash)
  II bob bo'yicha xulosa

III BOB (takomillashtirish yo'llari)
  3.1. (kichik mavzu — takliflar va tavsiyalar)
  3.2. (kichik mavzu — amalga oshirish mexanizmi)
  3.3. (kichik mavzu — kutilayotgan natijalar)
  III bob bo'yicha xulosa

UMUMIY XULOSA
FOYDALANILGAN ADABIYOTLAR RO'YXATI (kamida 15-20 ta manba)
GLOSSARIY (asosiy atamalar va ta'riflari)
ILOVALAR
""",

    "Magistrlik dissertatsiyasi": """
Magistrlik dissertatsiyasi uchun yuqori ilmiy darajadagi, TO'LIQ va JUDA BATAFSIL akademik
tuzilishda yozing (umumiy hajmi 125-130 betga mo'ljallangan, shuning uchun har bir kichik
bo'limni (masalan 1.1, 1.2, 1.3) alohida, chuqur, keng va ilmiy asoslangan tarzda yoriting):

MUNDARIJA (I, II, III bob asosida)

KIRISH
- Mavzuning dolzarbligi
- Ilmiy muammoning qo'yilishi
- Tadqiqotning maqsadi va vazifalari
- Tadqiqot ob'ekti, predmeti va metodologiyasi
- Ilmiy yangiligi
- Nazariy va amaliy ahamiyati

I BOB (ilmiy-nazariy tahlil)
  1.1. (kichik mavzu — mavzuning nazariy-metodologik asoslari)
  1.2. (kichik mavzu — mavjud tadqiqotlar va ilmiy qarashlar sharhi)
  1.3. (kichik mavzu — xorijiy va mahalliy tajriba tahlili)
  I bob bo'yicha xulosa

II BOB (metodologiya va empirik tahlil)
  2.1. (kichik mavzu — tadqiqot metodologiyasi)
  2.2. (kichik mavzu — ma'lumotlar to'plash va tahlil usullari)
  2.3. (kichik mavzu — empirik natijalarning dastlabki tahlili)
  II bob bo'yicha xulosa

III BOB (natijalar va tavsiyalar)
  3.1. (kichik mavzu — natijalarning ilmiy talqini)
  3.2. (kichik mavzu — amaliy tavsiyalar)
  3.3. (kichik mavzu — qo'llanish istiqbollari)
  III bob bo'yicha xulosa

UMUMIY XULOSA
FOYDALANILGAN ADABIYOTLAR RO'YXATI (125-130 ta manba, ko'pchiligi xalqaro nashrlardan)
GLOSSARIY (asosiy ilmiy atamalar va ta'riflari)
ILOVALAR
""",

    "PhD dissertatsiya": """
PhD (falsafa doktori) dissertatsiyasi uchun eng yuqori ilmiy standartlarga mos, TO'LIQ va
ENG BATAFSIL akademik tuzilishda yozing (umumiy hajmi 130-140 betga mo'ljallangan, shuning
uchun har bir kichik bo'limni chuqur ilmiy tahlil, gipoteza tekshiruvi va argumentatsiya
bilan yoriting):

MUNDARIJA (I, II, III, IV bob asosida)

KIRISH
- Muammoning qo'yilishi va dolzarbligi
- Tadqiqot muammosi va ilmiy gipoteza
- Tadqiqotning maqsadi va vazifalari
- Tadqiqot ob'ekti va predmeti
- Tadqiqot metodologiyasi
- Ilmiy yangiligi
- Himoyaga chiqariladigan asosiy holatlar
- Ishning nazariy va amaliy ahamiyati
- Natijalarning aprobatsiyasi (nashrlar, konferensiyalar haqida umumiy tavsif)

I BOB (adabiyotlar tahlili va nazariy-metodologik asos)
  1.1. (kichik mavzu — mavzu bo'yicha fundamental nazariyalar)
  1.2. (kichik mavzu — xalqaro tadqiqotlar sharhi)
  1.3. (kichik mavzu — nazariy asosning shakllanishi)
  I bob bo'yicha xulosa

II BOB (tadqiqot metodologiyasi)
  2.1. (kichik mavzu — tadqiqot dizayni)
  2.2. (kichik mavzu — ma'lumotlar to'plash usullari)
  2.3. (kichik mavzu — tahlil metodlari va vositalari)
  II bob bo'yicha xulosa

III BOB (tadqiqot natijalari)
  3.1. (kichik mavzu — asosiy empirik natijalar)
  3.2. (kichik mavzu — natijalarning statistik/sifat tahlili)
  3.3. (kichik mavzu — gipotezalarni tekshirish)
  3.4. (kichik mavzu — natijalarning ilmiy-nazariy talqini)
  3.5. (kichik mavzu — mavjud tadqiqotlar bilan qiyosiy tahlil)
  3.6. (kichik mavzu — cheklovlar va ularni bartaraf etish yo'llari)
  III bob bo'yicha xulosa

IV BOB (amaliy tavsiyalar va qo'llanish istiqbollari)
  4.1. (kichik mavzu — amaliy tavsiyalar)
  4.2. (kichik mavzu — qo'llanish mexanizmi)
  4.3. (kichik mavzu — kelgusi tadqiqotlar uchun yo'nalishlar)
  IV bob bo'yicha xulosa

UMUMIY XULOSA (asosiy ilmiy natijalarning yig'ma bayoni)
FOYDALANILGAN ADABIYOTLAR RO'YXATI (130-140 ta manba, asosan xalqaro nashrlardan)
GLOSSARIY (asosiy ilmiy atamalar va ta'riflari)
ILOVALAR
""",

    "Uslubiy qo'llanma": """
Uslubiy qo'llanma uchun quyidagi tuzilishda yozing:
1. Kirish so'zi (qo'llanmaning maqsadi va kimlar uchun mo'ljallanganligi)
2. Mundarija
3. Mavzular bo'yicha metodik ko'rsatmalar (har bir mavzu uchun: maqsad, tushuntirish, misol)
4. Amaliy mashg'ulotlar uchun tavsiyalar
5. Nazorat savollari va topshiriqlar
6. Baholash mezonlari
7. Tavsiya etilgan adabiyotlar ro'yxati
""",

    "O'quv qo'llanma": """
O'quv qo'llanma uchun bitta silabus yoki bir nechta katta mavzularni qamrab oluvchi,
TO'LIQ va JUDA BATAFSIL tuzilishda yozing (4-8 bobdan iborat, umumiy hajmi 170-240 betga
mo'ljallangan, shuning uchun har bir kichik bo'limni chuqur va misollar bilan yoriting):

MUQADDIMA (qo'llanmaning maqsadi va kimlar uchun mo'ljallanganligi)
MUNDARIJA (boblar bo'yicha)

Har bir BOB uchun (4-8 ta bob bo'lishi kerak):
  I-BOB (masalan)
    1.1. (kichik mavzu — nazariy material)
    1.2. (kichik mavzu — nazariy material, misollar bilan)
    1.3. (kichik mavzu — amaliy tahlil)
    Bob bo'yicha xulosa
    Bob bo'yicha nazorat savollari (5-10 ta)
    Bob bo'yicha nazorat testlari (5-10 ta, variantlar bilan)

(Shu tarzda barcha bobларни ketma-ket davom ettiring)

YAKUNIY XULOSA
GLOSSARIY (asosiy atamalar va ta'riflari)
FOYDALANILGAN VA TAVSIYA ETILGAN ADABIYOTLAR RO'YXATI
""",

    "Darslik": """
To'liq darslik uchun bitta fan dasturini to'liq qamrab oluvchi, ENG BATAFSIL va CHUQUR
tuzilishda yozing (6-14 bobdan iborat, umumiy hajmi 170-440 betga mo'ljallangan, shuning
uchun har bir kichik bo'limni maksimal darajada batafsil, misollar va amaliy masalalar
bilan yoriting):

MUQADDIMA (darslikning maqsadi, tuzilishi, fan dasturi bilan bog'liqligi, kimlar uchun
mo'ljallanganligi)
MUNDARIJA (boblar va kichik bo'limlar bo'yicha)

Har bir BOB uchun (6-14 ta bob bo'lishi kerak):
  I-BOB (masalan)
    1.1. (kichik mavzu — nazariy qism, chuqur va tushunarli tarzda)
    1.2. (kichik mavzu — nazariy qism, amaliy misollar bilan)
    1.3. (kichik mavzu — qo'shimcha chuqurlashtirilgan material)
    Bob bo'yicha xulosa
    Bob bo'yicha nazorat savollari (8-12 ta)
    Bob bo'yicha nazorat testlari (8-12 ta, variantlar bilan)

(Shu tarzda barcha boblarni ketma-ket davom ettiring)

UMUMIY XULOSA
GLOSSARIY (barcha atamalar va ta'riflari)
FOYDALANILGAN ADABIYOTLAR RO'YXATI
""",
}

DEFAULT_STRUCTURE = """
Mavzuni aniq, tushunarli va professional tarzda, kirish, asosiy qism va xulosadan
iborat tuzilishda yoriting.
"""


def get_structure_guide(context):
    """Bo'lim nomiga mos akademik tuzilma (andoza) qaytaradi."""
    return STRUCTURE_GUIDES.get(context, DEFAULT_STRUCTURE)


def _call_ai(system_prompt, prompt):
    """
    Avval OpenAI'ni sinaydi, ishlamasa Claude'ga, u ham ishlamasa
    Gemini'ga o'tadi. Uchalasi ham ishlamasa xatolik ko'taradi.
    """
    providers = [
        ("OpenAI", _try_openai),
        ("Claude", _try_claude),
        ("Gemini", _try_gemini),
    ]

    errors = []
    for name, func in providers:
        try:
            return func(system_prompt, prompt)
        except Exception as e:
            errors.append(f"{name}: {str(e)}")
            continue

    raise Exception("Barcha provayderlar ishlamadi:\n" + "\n".join(errors))


def get_ai_response(prompt, context="", language="uz"):
    """
    Qisqa hujjatlar uchun — bitta so'rov bilan to'liq javob oladi.
    """
    structure = get_structure_guide(context)
    lang_name = get_language_name(language)

    if context == "Maqola":
        same_as_ru = language == "ru"
        same_as_en = language == "en"
        extra_note = ""
        if same_as_ru or same_as_en:
            extra_note = (
                "\nEslatma: tanlangan til rus yoki ingliz tili bilan bir xil bo'lsa, "
                "takrorlanadigan tilni qayta yozmang, faqat 2 xil versiya bering."
            )
        system_prompt = (
            "Siz akademik AI yordamchisiz. Bo'lim: Ilmiy maqola (IMRAD talabi asosida).\n"
            f"Maqolaning asosiy mazmuni (Kirish, Materiallar va metodlar, Natijalar, "
            f"Muhokama, Xulosa qismlari) albatta {lang_name} tilida yozilishi kerak.\n"
            "Ammo SARLAVHA, ANNOTATSIYA va KALIT SO'ZLAR albatta UCH TILDA berilishi shart: "
            f"1) {lang_name} tilida, 2) rus tilida, 3) ingliz tilida." + extra_note +
            "\nProfessional, grammatik jihatdan to'g'ri va ilmiy uslubda yozing.\n\n"
            f"YOZISH TARTIBI (albatta shu tuzilishga qat'iy amal qiling):\n{structure}"
        )
    else:
        system_prompt = (
            f"Siz akademik AI yordamchisiz. Bo'lim: {context}. "
            f"MUHIM: Butun javobni albatta {lang_name} tilida yozing (sarlavhalar, "
            f"bo'lim nomlari va matnning barcha qismlari ham shu tilda bo'lsin).\n"
            "O'qituvchi va talabalar uchun professional, yaxshi tuzilgan va "
            "grammatik jihatdan to'g'ri javob bering.\n\n"
            f"YOZISH TARTIBI (albatta shu tuzilishga qat'iy amal qiling, lekin "
            f"matnning o'zini {lang_name} tilida yozing):\n{structure}"
        )

    try:
        return _call_ai(system_prompt, prompt)
    except Exception as e:
        return f"AI Xatolik yuz berdi. {str(e)}"


# ==========================================================
#  KATTA HAJMLI HUJJATLARNI BO'LIM-BO'LIM (BOB-BOB) YARATISH
# ==========================================================

def _chapter_sections(bob_nomi, sub_topics_count=3, has_tests=False):
    """Bitta bob uchun kichik bo'limlar ro'yxatini generatsiya qiladi."""
    sections = []
    for i in range(1, sub_topics_count + 1):
        sections.append(f"{bob_nomi} — {i}-band (kichik mavzu, chuqur va batafsil tahlil)")
    sections.append(f"{bob_nomi} bo'yicha xulosa")
    if has_tests:
        sections.append(f"{bob_nomi} bo'yicha nazorat savollari (8-12 ta)")
        sections.append(f"{bob_nomi} bo'yicha nazorat testlari (8-12 ta, variantlar bilan)")
    return sections


LONG_DOCUMENT_PLANS = {
    "Kurs ishi": (
        ["Kirish (dolzarblik, maqsad va vazifalar, tadqiqot ob'ekti/predmeti)"]
        + _chapter_sections("I BOB (nazariy qism)", 2)
        + _chapter_sections("II BOB (amaliy/tahliliy qism)", 2)
        + ["Umumiy xulosa va takliflar", "Foydalanilgan adabiyotlar ro'yxati (8-10 ta manba)"]
    ),
    "Bitiruv malakaviy ishi": (
        ["Kirish (dolzarblik, maqsad va vazifalar, tadqiqot ob'ekti/predmeti, ilmiy-amaliy ahamiyat)"]
        + _chapter_sections("I BOB (nazariy-metodologik asoslar)", 3)
        + _chapter_sections("II BOB (amaliy tahlil)", 3)
        + _chapter_sections("III BOB (takomillashtirish yo'llari)", 3)
        + ["Umumiy xulosa", "Foydalanilgan adabiyotlar ro'yxati (15-20 ta manba)", "Glossariy (asosiy atamalar)"]
    ),
    "Magistrlik dissertatsiyasi": (
        ["Kirish (dolzarblik, ilmiy muammo, maqsad-vazifalar, ob'ekt/predmet, metodologiya, ilmiy yangilik, ahamiyat)"]
        + _chapter_sections("I BOB (ilmiy-nazariy tahlil)", 3)
        + _chapter_sections("II BOB (metodologiya va empirik tahlil)", 3)
        + _chapter_sections("III BOB (natijalar va tavsiyalar)", 3)
        + ["Umumiy xulosa", "Foydalanilgan adabiyotlar ro'yxati (25-30 ta manba)", "Glossariy (asosiy ilmiy atamalar)"]
    ),
    "PhD dissertatsiya": (
        ["Kirish (muammo, gipoteza, maqsad-vazifalar, ob'ekt/predmet, metodologiya, ilmiy yangilik, "
         "himoyaga chiqariladigan holatlar, ahamiyat, aprobatsiya)"]
        + _chapter_sections("I BOB (adabiyotlar tahlili va nazariy-metodologik asos)", 3)
        + _chapter_sections("II BOB (tadqiqot metodologiyasi)", 3)
        + _chapter_sections("III BOB (tadqiqot natijalari)", 6)
        + _chapter_sections("IV BOB (amaliy tavsiyalar va istiqbollar)", 3)
        + ["Umumiy xulosa", "Foydalanilgan adabiyotlar ro'yxati (50+ manba)", "Glossariy (asosiy ilmiy atamalar)"]
    ),
    "O'quv qo'llanma": (
        ["Muqaddima (qo'llanmaning maqsadi va kimlar uchun mo'ljallanganligi)"]
        + _chapter_sections("I BOB", 3, has_tests=True)
        + _chapter_sections("II BOB", 3, has_tests=True)
        + _chapter_sections("III BOB", 3, has_tests=True)
        + _chapter_sections("IV BOB", 3, has_tests=True)
        + ["Yakuniy xulosa", "Glossariy (asosiy atamalar)", "Foydalanilgan va tavsiya etilgan adabiyotlar ro'yxati"]
    ),
    "Darslik": (
        ["Muqaddima (darslikning maqsadi, tuzilishi, fan dasturi bilan bog'liqligi)"]
        + _chapter_sections("I BOB", 3, has_tests=True)
        + _chapter_sections("II BOB", 3, has_tests=True)
        + _chapter_sections("III BOB", 3, has_tests=True)
        + _chapter_sections("IV BOB", 3, has_tests=True)
        + _chapter_sections("V BOB", 3, has_tests=True)
        + _chapter_sections("VI BOB", 3, has_tests=True)
        + ["Umumiy xulosa", "Glossariy (barcha atamalar)", "Foydalanilgan adabiyotlar ro'yxati"]
    ),
}


def is_long_document(category):
    """Bu bo'lim bo'lim-bo'lim generatsiya qilinishi kerakligini tekshiradi."""
    return category in LONG_DOCUMENT_PLANS


def generate_long_document(topic, category, progress_callback=None, language="uz"):
    """
    Katta hajmli hujjatlarni (Kurs ishi, BMI, Magistrlik, PhD, O'quv qo'llanma, Darslik)
    bo'lim-bo'lim generatsiya qiladi va birlashtirib qaytaradi.
    """
    plan = LONG_DOCUMENT_PLANS.get(category)
    if not plan:
        return get_ai_response(topic, context=category, language=language)

    lang_name = get_language_name(language)

    base_system_prompt = (
        f"Siz akademik AI yordamchisiz. Mavzu: \"{topic}\". Hujjat turi: {category}. "
        f"MUHIM: Butun javobni albatta {lang_name} tilida yozing. "
        "Professional, ilmiy uslubda va grammatik jihatdan to'g'ri yozing. "
        "FAQAT so'ralgan bo'limga oid matnni yozing — boshqa bo'limlarni takrorlamang, "
        "umumiy kirish so'zi yoki xulosa qo'shmang, faqat so'ralgan qism matnini bering."
    )

    total = len(plan)
    parts = []

    for idx, section_title in enumerate(plan, start=1):
        section_prompt = (
            f"Mavzu: \"{topic}\"\n\n"
            f"Quyidagi bo'lim uchun batafsil, professional va ilmiy asoslangan matn yozing:\n"
            f"\"{section_title}\"\n\n"
            "Matn kamida 300-500 so'zdan iborat, chuqur va mazmunli bo'lsin."
        )
        try:
            text = _call_ai(base_system_prompt, section_prompt)
        except Exception as e:
            text = f"[Ushbu bo'limni yaratishda xatolik yuz berdi: {e}]"

        parts.append(f"\n\n{section_title.upper()}\n\n{text}")

        if progress_callback:
            try:
                progress_callback(idx, total, section_title)
            except Exception:
                pass

    return "".join(parts)


# ==========================================================
#  WORD FAYL YARATISH
# ==========================================================

def create_word(title, content):
    doc = Document()
    doc.add_heading(title, 0)
    for line in content.split('\n'):
        if line.strip():
            doc.add_paragraph(line)
    file_path = "Akademik_Ish.docx"
    doc.save(file_path)
    return file_path


# ==========================================================
#  PEXELS ORQALI MAVZUGA MOS RASM TOPISH
# ==========================================================

def _search_pexels_image(query, save_path="temp_slide_image.jpg"):
    """
    Pexels'dan mavzuga mos rasm qidiradi va yuklab oladi.
    """
    if not PEXELS_API_KEY:
        return None
    try:
        headers = {"Authorization": PEXELS_API_KEY}
        params = {"query": query, "per_page": 1, "orientation": "landscape"}
        resp = requests.get(
            "https://api.pexels.com/v1/search",
            headers=headers, params=params, timeout=10
        )
        data = resp.json()
        photos = data.get("photos", [])
        if not photos:
            return None
        image_url = photos[0]["src"]["large"]
        img_resp = requests.get(image_url, timeout=10)
        with open(save_path, "wb") as f:
            f.write(img_resp.content)
        return save_path
    except Exception:
        return None


# ==========================================================
#  POWERPOINT YARATISH — professional dizayn, mavzuga aynan mos
#  rasmlar (avtomatik ingliz tiliga tarjima qilib qidiriladi)
#  va jadval ko'rinishidagi ma'lumotlar uchun haqiqiy pptx jadvali
# ==========================================================

_TRANSLATION_CACHE = {}


def _translate_to_english(text):
    """Qisqa matnni (slayd sarlavhasi) rasm qidiruvi uchun ingliz tiliga
    tarjima qiladi. Xatolik bo'lsa asl matnni qaytaradi."""
    if not text:
        return text
    if text in _TRANSLATION_CACHE:
        return _TRANSLATION_CACHE[text]
    try:
        result = _call_ai(
            "You are a translator. Translate the given short phrase into English. "
            "Respond with ONLY the translated phrase (2-4 words), no quotes, no "
            "explanation, no extra text.",
            text
        )
        cleaned = result.strip().strip('"').strip("'").split("\n")[0][:60]
        _TRANSLATION_CACHE[text] = cleaned if cleaned else text
        return _TRANSLATION_CACHE[text]
    except Exception:
        return text


def create_pptx(title, content):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    ACCENT_COLOR = RGBColor(0x1F, 0x4E, 0x79)   # to'q ko'k
    ACCENT_LIGHT = RGBColor(0xDD, 0xE9, 0xF5)   # och ko'k (jadval sarlavhasi uchun)
    TEXT_COLOR = RGBColor(0x2B, 0x2B, 0x2B)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # bo'sh maket — to'liq nazorat uchun

    paragraphs = [p for p in content.split('\n\n') if len(p.strip()) > 10]

    for p in paragraphs[:25]:  # Maksimal 25 ta slayd
        lines = [ln for ln in p.strip().split('\n') if ln.strip()]
        if not lines:
            continue

        first_line = lines[0].strip()
        if len(first_line) <= 90 and len(lines) > 1:
            slide_title = first_line.rstrip(':').strip()
            body_lines = lines[1:]
        else:
            slide_title = title
            body_lines = lines

        slide = prs.slides.add_slide(blank_layout)

        # --- Fon ---
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE

        # --- Yuqori aksent chizig'i ---
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.18)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT_COLOR
        bar.line.fill.background()
        bar.shadow.inherit = False

        # --- Sarlavha ---
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.9))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.text = slide_title
        title_tf.paragraphs[0].font.size = Pt(26)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = ACCENT_COLOR

        # --- Jadval va oddiy matn qatorlarini ajratish ---
        table_lines = [ln for ln in body_lines if '|' in ln]
        text_lines = [ln for ln in body_lines if '|' not in ln]

        content_top = 1.4
        has_table = bool(table_lines)

        # --- Matn qismi (bullet ko'rinishida) ---
        if text_lines:
            text_width = 5.6 if not has_table else 9.0
            body_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(content_top), Inches(text_width), Inches(4.8)
            )
            body_tf = body_box.text_frame
            body_tf.word_wrap = True
            for i, line in enumerate(text_lines):
                clean_line = line.lstrip('-•').strip()
                para = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
                para.text = f"•  {clean_line}"
                para.font.size = Pt(18)
                para.font.color.rgb = TEXT_COLOR
                para.space_after = Pt(10)
            content_top += 0.3 + 0.5 * len(text_lines)

        # --- Jadval qismi (haqiqiy PowerPoint jadvali) ---
        if has_table:
            rows_data = [[c.strip() for c in ln.split('|')] for ln in table_lines]
            n_rows = len(rows_data)
            n_cols = max(len(r) for r in rows_data)
            table_top = Inches(content_top) if text_lines else Inches(1.6)
            table_height = min(0.5 * n_rows, 4.5)
            table_shape = slide.shapes.add_table(
                n_rows, n_cols, Inches(0.5), table_top, Inches(9.0), Inches(table_height)
            )
            table = table_shape.table
            for r_idx, row in enumerate(rows_data):
                for c_idx in range(n_cols):
                    cell = table.cell(r_idx, c_idx)
                    cell.text = row[c_idx] if c_idx < len(row) else ""
                    for para in cell.text_frame.paragraphs:
                        para.font.size = Pt(13)
                        if r_idx == 0:
                            para.font.bold = True
                    if r_idx == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = ACCENT_LIGHT

        # --- Mavzuga aynan mos rasm (ingliz tiliga tarjima qilib qidiriladi) ---
        if not has_table:
            search_term = slide_title if slide_title != title else title
            english_query = _translate_to_english(search_term)
            image_path = _search_pexels_image(english_query)
            if image_path and os.path.exists(image_path):
                try:
                    slide.shapes.add_picture(
                        image_path, Inches(6.4), Inches(1.6), width=Inches(3.1)
                    )
                except Exception:
                    pass
                finally:
                    if os.path.exists(image_path):
                        os.remove(image_path)

    file_path = "Taqdimot.pptx"
    prs.save(file_path)
    return file_path
